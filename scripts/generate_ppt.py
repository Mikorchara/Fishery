"""
毕业答辩 PPT 生成器 — 智慧渔业水下协同控制系统
运行: python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# 16:9 宽屏
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# 颜色主题 — 深蓝科技风
# ============================================================
CLR_DARK = RGBColor(0x0B, 0x1D, 0x3A)        # 深蓝背景
CLR_ACCENT = RGBColor(0x00, 0xB4, 0xD8)       # 青蓝强调
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_LIGHT = RGBColor(0xE0, 0xF2, 0xFE)
CLR_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
CLR_GREEN = RGBColor(0x10, 0xB9, 0x81)
CLR_GRAY = RGBColor(0x94, 0xA3, 0xB8)
CLR_SUBTITLE = RGBColor(0x7D, 0xC4, 0xE0)


def dark_slide():
    """创建深蓝背景页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = CLR_DARK
    return slide


def add_title(slide, text, left=1, top=0.6, width=11.3, height=1.2, font_size=40):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CLR_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return tf


def add_subtitle(slide, text, left=1, top=2.0, width=11.3, font_size=20):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CLR_SUBTITLE
    p.alignment = PP_ALIGN.LEFT
    return tf


def add_body(slide, text_lines, left=1, top=2.8, width=11.3, font_size=18, line_spacing=1.6):
    """添加多行正文，支持 • 前缀"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = CLR_LIGHT
        p.space_after = Pt(font_size * (line_spacing - 1))
    return tf


def add_accent_line(slide, left=1, top=1.85, width=3):
    """青蓝色装饰线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_ACCENT
    shape.line.fill.background()
    return shape


def add_page_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.0), Inches(1), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}"
    p.font.size = Pt(12)
    p.font.color.rgb = CLR_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_card(slide, title, lines, left, top, width=3, height=2.5, accent_color=CLR_ACCENT):
    """卡片模块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x12, 0x2D, 0x52)
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    p.font.bold = True
    # 内容
    txBox2 = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.65), Inches(width - 0.4), Inches(height - 0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = CLR_LIGHT
        p.space_after = Pt(4)
    return shape


# ============================================================
# Slide 1: 封面
# ============================================================
slide = dark_slide()
add_title(slide, "智慧渔业水下协同控制系统", top=1.5, font_size=48)
add_accent_line(slide, top=2.8, width=4)
add_subtitle(slide, "基于深度学习的水下目标检测与图像增强", top=3.1, font_size=22)
add_body(slide, ["毕业设计答辩", "导 师：______       答辩人：______", "计算机科学与技术 / 人工智能方向", "2026年4月"], top=4.0, font_size=16)
add_page_number(slide, 1)

# ============================================================
# Slide 2: 目录
# ============================================================
slide = dark_slide()
add_title(slide, "目  录")
add_accent_line(slide)
toc = [
    "01   项目背景与意义",
    "02   系统总体架构",
    "03   核心技术 — 水下图像增强 (WWE-UIE)",
    "04   核心技术 — 目标检测与实例分割 (YOLOv8 + SAM2)",
    "05   核心技术 — AI 智慧养殖顾问 (LLM)",
    "06   系统实现与 Web 交互界面",
    "07   关键难点与解决方案",
    "08   总结与展望",
]
add_body(slide, toc, top=2.6, font_size=20, line_spacing=1.8)
add_page_number(slide, 2)

# ============================================================
# Slide 3: 项目背景与意义
# ============================================================
slide = dark_slide()
add_title(slide, "01  项目背景与意义")
add_accent_line(slide)
add_body(slide, [
    "• 水产养殖是全球重要的蛋白质来源，但传统养殖依赖人工巡检，效率低、成本高",
    "• 水下环境光照不足、色彩偏蓝/绿、对比度低，肉眼难以准确判断鱼群状态",
    "• 市场上缺乏一套集「画面增强 → 自动检测 → 智能诊断」于一体的完整系统",
    "",
    "🎯  本项目目标：",
    "• 构建一个端到端的智慧渔业监控平台，实现水下视频实时采集、画质修复、",
    "  鱼群自动检测/分割，以及基于 LLM 的养殖环境智能分析",
    "• 将深度学习技术（YOLOv8、SAM2、WWE-UIE）从论文落地到工程实践",
], top=2.4, font_size=17)
add_page_number(slide, 3)

# ============================================================
# Slide 4: 系统总体架构
# ============================================================
slide = dark_slide()
add_title(slide, "02  系统总体架构")
add_accent_line(slide)
# 四层卡片
add_card(slide, "🖥 用户交互层", ["• Web 实时视频预览", "• AI 功能开关控制", "• 传感器数据面板", "• 智慧对话交互"], 0.6, 2.6, 2.9, 2.3, CLR_ACCENT)
add_card(slide, "⚙ 应用调度层", ["• Flask Web Server", "• MJPEG 视频流推送", "• RESTful API 路由", "• 全局状态管理"], 3.7, 2.6, 2.9, 2.3, CLR_GREEN)
add_card(slide, "🧠 核心处理层", ["• YOLOv8 目标检测", "• SAM2 实例分割", "• WWE-UIE 图像增强", "• NVIDIA LLM 顾问"], 6.8, 2.6, 2.9, 2.3, CLR_GOLD)
add_card(slide, "💾 数据支撑层", ["• 模型权重 (.pt/.pth)", "• RTSP 视频流接入", "• IoT 传感器数据", "• 系统配置中心"], 9.9, 2.6, 2.9, 2.3, CLR_GRAY)
# 数据流箭头描述
add_body(slide, [
    "数据流向： RTSP 视频源 → 双线程采集 → 图像增强 (可选) → AI 检测/分割 (可选) → MJPEG 推送 → Web 浏览器",
], top=5.3, font_size=14)
add_page_number(slide, 4)

# ============================================================
# Slide 5: 视频采集模块
# ============================================================
slide = dark_slide()
add_title(slide, "02a  视频采集 — 双线程异步架构")
add_accent_line(slide)
add_body(slide, [
    "▎ 设计动机",
    "• OpenCV 读取 RTSP 流时网络抖动会导致主线程阻塞，画面卡顿",
    "",
    "▎ 解决方案：双线程分离",
    "• 后台线程持续抓帧存到内存缓冲区，主线程异步取帧，互不阻塞",
    "• 强制 TCP 传输（非 UDP），避免花屏和丢包",
    "",
    "▎ 关键实现",
    "• VideoCaptureThreading 类：daemon 线程 + 0.01s 休息间隔",
    "• read() 接口无锁返回最新帧，实现「零等待」消费",
], top=2.4, font_size=16)
add_page_number(slide, 5)

# ============================================================
# Slide 6: WWE-UIE
# ============================================================
slide = dark_slide()
add_title(slide, "03  水下图像增强 — WWE-UIE 算法")
add_accent_line(slide)
add_body(slide, [
    "▎ 为什么需要水下增强？",
    "• 水体对红光吸收强烈，水下画面偏蓝/绿",
    "• 悬浮颗粒导致雾化、对比度下降 → 直接影响 YOLO 检测准确率",
    "",
    "▎ WWE-UIE 网络结构",
    "• 基于 CNN，包含白平衡模块 (White Balance) + 色彩校正 + 对比度增强",
    "• 输入 RGB 三通道，输出恢复后的自然色彩图像",
    "• 在 UIEB 等水下数据集上训练，支持 CUDA 实时推理",
    "",
    "▎ 工作流程",
    "BGR帧 → RGB归一化 → Tensor(1,3,H,W) → 模型推理 → 后处理 → BGR输出 → 送入检测模块",
], top=2.4, font_size=15)
add_page_number(slide, 6)

# ============================================================
# Slide 7: YOLOv8 + SAM2
# ============================================================
slide = dark_slide()
add_title(slide, "04  目标检测 + 实例分割 — YOLOv8 & SAM2")
add_accent_line(slide)
add_body(slide, [
    "▎ YOLOv8 负责快速定位",
    "• 使用 Ultralytics 框架，加载 fish_detect.pt（专门针对鱼群场景训练）",
    "• 可调节：置信度阈值(0.15)、IOU(0.3)、推理尺寸(640)、最大检测数(100)",
    "• 支持热切换模型权重，无需重启服务",
    "",
    "▎ SAM2 负责像素级分割（可选开启）",
    "• 在 YOLO 检测框基础上进一步提取精确轮廓",
    "• 针对鳗鱼等长条形目标采用 Point Prompt 策略（质心点 + 标签=1）",
    "• 支持三种提示模式: box / point / hybrid，可配置切换",
    "",
    "▎ 工作流程",
    "增强帧 → YOLO预测 → (不开分割)直接绘制检测框 → 输出",
    "                        → (开分割)提取boxes → SAM2 精细化遮罩 → 输出",
], top=2.4, font_size=14)
add_page_number(slide, 7)

# ============================================================
# Slide 8: LLM Advisor
# ============================================================
slide = dark_slide()
add_title(slide, "05  AI 智慧养殖顾问 — LLM 集成")
add_accent_line(slide)
add_body(slide, [
    "▎ 集成方式",
    "• 通过 OpenAI SDK 调用 NVIDIA API 端点（MiniMax-M2.7 模型）",
    "• System Prompt 设定为「资深水产养殖专家」角色",
    "",
    "▎ 两大功能",
    "1. 环境诊断报告：读取 IoT 传感器数据（水温/pH/溶解氧），自动生成养殖评估与建议",
    "2. 自由对话：用户在网页端输入问题（如「水有点浑浊怎么办」），AI 结合实时数据给出建议",
    "",
    "▎ 设计考量",
    "• 使用非流式响应 (stream=False)，确保 Flask 同步路由兼容",
    "• 上下文注入：每次请求都附上当前传感器读数作为环境参考",
], top=2.4, font_size=16)
add_page_number(slide, 8)

# ============================================================
# Slide 9: Web 界面
# ============================================================
slide = dark_slide()
add_title(slide, "06  Web 交互界面")
add_accent_line(slide)
add_body(slide, [
    "▎ 技术栈：HTML5 + CSS3 + Vanilla JS（无框架依赖，轻量部署）",
    "",
    "▎ 界面分区",
    "• 左侧主区域：实时视频画面（MJPEG 流，<img> 标签直接渲染）",
    "• 右侧控制面板：",
    "    - IoT 传感器数据卡片（水温/pH/溶解氧，每 2s 自动刷新）",
    "    - AI 视觉中枢（开关检测/增强/分割 + 模型选择器）",
    "    - 智慧对话窗口（Markdown 渲染，支持上下文对话）",
    "",
    "▎ 前后端交互",
    "• 开关按钮 → POST 请求 → Flask 路由 → 修改全局 system_state",
    "• 传感器轮询 → GET /get_sensor 每 2s",
    "• AI 对话 → POST /chat_ai，返回 Markdown，前端 marked.js 渲染",
], top=2.4, font_size=15)
add_page_number(slide, 9)

# ============================================================
# Slide 10: 难点与方案
# ============================================================
slide = dark_slide()
add_title(slide, "07  关键难点与解决方案")
add_accent_line(slide)
add_body(slide, [
    "🔴 难点 1：视频卡顿",
    "    原因：AI 推理耗时阻塞视频采集主循环",
    "    ✅  方案：双线程架构 — 采集线程独立运行，推理线程异步消费",
    "",
    "🔴 难点 2：水下画面质量差，检测率低",
    "    原因：水体偏色、低对比度导致 YOLO 提取不到鱼群特征",
    "    ✅  方案：在 YOLO 上游加入 WWE-UIE 增强模块，先修复画质再检测",
    "",
    "🔴 难点 3：鳗鱼型目标分割不完整",
    "    原因：检测框过大/长条形物体框内包含大量背景 → SAM 无焦点",
    "    ✅  方案：切换到 Point Prompt 模式，用质心点引导 SAM 精确分割",
    "",
    "🔴 难点 4：多模型管理与内存占用",
    "    原因：SAM2 模型较大（~65MB），常驻内存浪费资源",
    "    ✅  方案：按需加载策略，默认只加载 YOLO，分割功能可选开启",
], top=2.3, font_size=13.5)
add_page_number(slide, 11)

# ============================================================
# Slide 11: 成果展示
# ============================================================
slide = dark_slide()
add_title(slide, "07b  系统运行效果")
add_accent_line(slide)
add_card(slide, "📊 检测性能", ["• 鱼群检测置信度≥0.15", "• 单帧推理 < 50ms (GPU)", "• 支持 100+ 目标同时检出", "• FPS 实时显示"], 0.6, 2.5, 3.6, 2.2, CLR_GREEN)
add_card(slide, "🎨 增强效果", ["• 色彩偏差明显修复", "• 对比度显著提升", "• 白平衡自适应调整", "• 不影响检测帧率"], 4.7, 2.5, 3.6, 2.2, CLR_ACCENT)
add_card(slide, "🤖 LLM 诊断", ["• 环境指标综合分析", "• 养殖风险预警", "• Markdown 格式化输出", "• 支持自由问答"], 8.8, 2.5, 3.6, 2.2, CLR_GOLD)
add_body(slide, ["说明：受限于答辩环境，建议准备一段实际运行的录屏视频在 Slide 12 播放"], top=5.3, font_size=14)
add_page_number(slide, 12)

# ============================================================
# Slide 12: 总结
# ============================================================
slide = dark_slide()
add_title(slide, "08  总结与展望")
add_accent_line(slide)
add_body(slide, [
    "✅  已完成工作",
    "• 构建了完整的「采集→增强→检测→分割→诊断→展示」端到端系统",
    "• 成功集成 YOLOv8 + SAM2 + WWE-UIE 三种深度学习模型",
    "• 实现了基于 LLM 的智能水产养殖顾问",
    "• 开发了直观的 Web 实时监控界面",
    "",
    "🔮  未来展望",
    "• 增加目标追踪（ByteTrack / BoT-SORT），实现鱼群轨迹分析",
    "• 引入多路摄像头同时监控，扩展为分布式系统",
    "• 接入真实边缘计算设备（Jetson / 树莓派）进行部署测试",
    "• 导出检测数据到数据库，支持历史数据回溯与统计分析",
], top=2.4, font_size=16)
add_page_number(slide, 13)

# ============================================================
# Slide 13: 致谢
# ============================================================
slide = dark_slide()
add_title(slide, "致  谢", font_size=48)
add_accent_line(slide, top=2.8, width=3)
add_body(slide, [
    "感谢各位老师抽出宝贵时间审阅本文并给予指导",
    "", "", "",
    "Q & A",
    "欢迎提问 🙏",
], top=3.5, font_size=20)
add_page_number(slide, 14)

# ============================================================
# 保存
# ============================================================
output_path = "毕业答辩_智慧渔业水下协同控制系统.pptx"
prs.save(output_path)
print(f"PPT generated: {output_path}")
print(f"Total slides: {len(prs.slides)}, 16:9 format")
